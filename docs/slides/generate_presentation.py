#!/usr/bin/env python3
"""
Générateur de présentation PowerPoint pour le 11e Logi Template (& Todolist) Builder.
Exécuter: python generate_presentation.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
BACKGROUND_IMAGE = os.path.join(SCRIPT_DIR, 'fond_pptx.png')

# Layout (en inches) — 16:9 13.333×7.5
LAYOUT = {
    'left_margin': 0.5,
    'right_margin': 0.5,
    'content_width': 12.333,
    'title_top': 0.35,
    'content_top': 1.5,
    'footer_top': 6.5,
    'center_x': 6.666,
}

# Couleurs du thème — Foxhole / 11eRC-FL
COLORS = {
    'primary':          RGBColor(88,  101, 242),   # Discord Blurple
    'secondary':        RGBColor(47,   49,  54),   # Discord Dark
    'accent':           RGBColor(87,  242, 135),   # Vert Discord
    'accent_gold':      RGBColor(255, 193,   7),   # Or / Highlights
    'accent_orange':    RGBColor(250, 130,  49),   # Orange Foxhole
    'text_main':        RGBColor(220, 220, 230),   # Gris clair
    'text_light':       RGBColor(255, 255, 255),   # Blanc pur
    'text_muted':       RGBColor(160, 160, 175),   # Secondaire
    'box_bg':           RGBColor( 45,  50,  65),   # Fond encadrés
    'box_border':       RGBColor( 88, 101, 242),   # Bordure encadrés
    'placeholder_bg':   RGBColor( 60,  65,  80),
    'placeholder_border': RGBColor(100, 110, 140),
    'tag_bg':           RGBColor( 30,  35,  48),   # Fond tags foncé
    'separator':        RGBColor( 88, 101, 242),   # Ligne déco
}


# ─────────────────────────────────────────────
#  Helpers génériques
# ─────────────────────────────────────────────

def add_background(slide, prs):
    """Ajoute l'image de fond (si disponible) en arrière-plan."""
    if os.path.exists(BACKGROUND_IMAGE):
        bg = slide.shapes.add_picture(
            BACKGROUND_IMAGE,
            Inches(0), Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)
    else:
        print(f"⚠️  Image de fond non trouvée : {BACKGROUND_IMAGE}")


def add_title(slide, text, top=None, font_size=40, bold=True, color=None):
    """Titre principal (décalé après éventuel logo)."""
    if top is None:
        top = LAYOUT['title_top']
    if color is None:
        color = COLORS['text_light']
    box = slide.shapes.add_textbox(
        Inches(LAYOUT['left_margin']), Inches(top),
        Inches(LAYOUT['content_width']), Inches(0.9),
    )
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT
    return box


def add_separator(slide, top, width=None):
    """Ligne de séparation horizontale colorée."""
    if width is None:
        width = LAYOUT['content_width']
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(LAYOUT['left_margin']), Inches(top),
        Inches(width), Inches(0.04),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['separator']
    line.line.fill.background()
    return line


def add_bullet_text(slide, bullets, left=None, top=None, width=None, font_size=16,
                    line_spacing=Pt(10)):
    """Liste à puces avec support de niveaux (indentation "  ") et couleurs spéciales."""
    if left is None:
        left = LAYOUT['left_margin']
    if top is None:
        top = LAYOUT['content_top']
    if width is None:
        width = LAYOUT['content_width']

    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(5.5),
    )
    tf = box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = line_spacing

        if not bullet:
            p.text = ''
            continue

        if bullet.startswith('    '):          # Niveau 2
            p.text = '      ◦  ' + bullet.strip()
            p.font.size = Pt(font_size - 2)
            p.level = 2
            p.font.color.rgb = COLORS['text_muted']
        elif bullet.startswith('  '):          # Niveau 1
            p.text = '    •  ' + bullet.strip()
            p.font.size = Pt(font_size - 1)
            p.level = 1
            p.font.color.rgb = COLORS['text_main']
        elif (bullet.startswith('✨') or bullet.startswith('⚠️')
              or bullet.startswith('💡')):
            p.text = bullet
            p.font.size = Pt(font_size)
            p.font.color.rgb = COLORS['accent_gold']
        elif bullet.startswith('🔑') or bullet.startswith('🎯'):
            p.text = bullet
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = COLORS['accent']
        else:
            p.text = '•  ' + bullet if bullet else ''
            p.font.size = Pt(font_size)
            p.font.color.rgb = COLORS['text_main']

    return box


def add_info_box(slide, text, left=None, top=1.5, width=None, height=0.9,
                 bg=None, border=None, font_size=17, bold=False, center=True):
    """Encadré coloré (highlight, note, etc.)."""
    if left is None:
        left = LAYOUT['left_margin']
    if width is None:
        width = LAYOUT['content_width']
    if bg is None:
        bg = COLORS['box_bg']
    if border is None:
        border = COLORS['box_border']

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = border
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = COLORS['accent_gold']
    p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
    tf.anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_placeholder_box(slide, text, left=7.0, top=1.5, width=5.8, height=4.7):
    """Placeholder pour capture d'écran (cadre pointillé)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['placeholder_bg']
    shape.line.color.rgb = COLORS['placeholder_border']
    shape.line.width = Pt(2)
    shape.line.dash_style = 2

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f'📸  {text}'
    p.font.size = Pt(13)
    p.font.color.rgb = COLORS['text_muted']
    p.alignment = PP_ALIGN.CENTER
    tf.anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_tag(slide, text, left, top, width=1.6, height=0.4,
            bg=None, border=None, font_size=12):
    """Petit badge/tag coloré."""
    if bg is None:
        bg = COLORS['primary']
    if border is None:
        border = COLORS['primary']
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = border
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.CENTER
    tf.anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_footer(slide, text):
    """Texte de bas de slide."""
    box = slide.shapes.add_textbox(
        Inches(LAYOUT['left_margin']), Inches(LAYOUT['footer_top']),
        Inches(LAYOUT['content_width']), Inches(0.6),
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = COLORS['text_muted']
    p.alignment = PP_ALIGN.LEFT
    return box


def add_step_box(slide, number, text, left, top, width=2.8, height=0.75):
    """Étape numérotée dans un flow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary']
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f'{number}. {text}'
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.CENTER
    tf.anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_arrow(slide, left, top, direction='right'):
    """Flèche directionnelle."""
    shape_type = MSO_SHAPE.RIGHT_ARROW if direction == 'right' else MSO_SHAPE.DOWN_ARROW
    w, h = (0.4, 0.28) if direction == 'right' else (0.28, 0.4)
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['accent']
    shape.line.fill.background()
    return shape


# ─────────────────────────────────────────────
#  SLIDES
# ─────────────────────────────────────────────

def create_title_slide(prs):
    """Slide 1 — Page de titre."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Titre principal
    title = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(2.0))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = '11e Logi Template'
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.CENTER

    # Sous-titre
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12.333), Inches(0.9))
    tf2 = sub.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = '(& Todolist) Builder'
    p2.font.size = Pt(36)
    p2.font.color.rgb = COLORS['accent_orange']
    p2.alignment = PP_ALIGN.CENTER

    # Séparateur
    add_separator(slide, 4.75, width=12.333)

    # Tagline
    tag = slide.shapes.add_textbox(Inches(0.5), Inches(4.95), Inches(12.333), Inches(0.7))
    tf3 = tag.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = 'Outil de logistique Foxhole pour la 11eRC-FL'
    p3.font.size = Pt(22)
    p3.font.italic = True
    p3.font.color.rgb = COLORS['text_muted']
    p3.alignment = PP_ALIGN.CENTER

    # Badges
    add_tag(slide, '🎨 Templates PNG', 3.2, 5.85, width=2.4, height=0.45, font_size=13)
    add_tag(slide, '📋 MPF Todolists', 5.9, 5.85, width=2.4, height=0.45,
            bg=COLORS['accent_orange'], border=COLORS['accent_orange'], font_size=13)
    add_tag(slide, '💾 Sauvegarde JSON', 8.6, 5.85, width=2.6, height=0.45,
            bg=RGBColor(60, 120, 80), border=RGBColor(60, 120, 80), font_size=13)


def create_overview_slide(prs):
    """Slide 2 — C'est quoi ?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "🔭  C'est quoi le 11e Template Builder ?")
    add_separator(slide, 1.2)

    bullets = [
        'Application web TypeScript pour la logistique Foxhole',
        'Conçue spécifiquement pour la 11eRC-FL',
        '',
        '🎨  Mode Template :',
        '  Créer des visuels de stockpile (canvas libre, sections)',
        '  Exporter en PNG 1920×1080 pour Discord / forum',
        '',
        '📋  Mode MPF Todolist :',
        '  Construire des listes de commande avec coûts MPF calculés',
        '  Exporter au format texte Discord (lettres indicatrices 🇦 🇧 🇨…)',
        '',
        '💾  Sauvegarde locale automatique (localStorage) + JSON import/export',
    ]
    add_bullet_text(slide, bullets, font_size=17, top=1.45)

    add_info_box(
        slide,
        '🌐  Accessible dans le navigateur — aucun logiciel à installer',
        top=6.35, height=0.75,
        bg=COLORS['box_bg'], border=COLORS['accent'],
        font_size=16, bold=True,
    )


def create_modes_slide(prs):
    """Slide 3 — Vue d'ensemble des deux modes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, '⚡  Deux modes, un seul outil')
    add_separator(slide, 1.2)

    # Colonne Template
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.45), Inches(5.8), Inches(4.7),
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = COLORS['box_bg']
    left_box.line.color.rgb = COLORS['primary']
    left_box.line.width = Pt(2)

    tl = left_box.text_frame
    tl.word_wrap = True
    tl.margin_left = Pt(14)
    tl.margin_top = Pt(12)
    for i, line in enumerate([
        '🎨 Mode Template',
        '',
        '• Canvas libre 1920×1080',
        '• Sections déplaçables / redimensionnables',
        '• Sidebar avec tous les items Foxhole',
        '• Drag & drop vers les sections',
        '• Tailles d\'icônes : S / M / L / XL / XXL',
        '• Fond de carte personnalisable',
        '• Export PNG haute qualité',
    ]):
        p = tl.paragraphs[0] if i == 0 else tl.add_paragraph()
        p.text = line
        p.font.size = Pt(14 if i > 0 else 18)
        p.font.bold = (i == 0)
        p.font.color.rgb = COLORS['text_light'] if i == 0 else COLORS['text_main']
        p.space_after = Pt(6)

    # Colonne TodoList
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.9), Inches(1.45), Inches(5.9), Inches(4.7),
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = COLORS['box_bg']
    right_box.line.color.rgb = COLORS['accent_orange']
    right_box.line.width = Pt(2)

    tr = right_box.text_frame
    tr.word_wrap = True
    tr.margin_left = Pt(14)
    tr.margin_top = Pt(12)
    for i, line in enumerate([
        '📋 Mode MPF Todolist',
        '',
        '• Catalogue complet items Foxhole',
        '• Filtre par catégorie + recherche',
        '• Coûts MPF calculés automatiquement',
        '  (Bmats / Rmats / Emats / HEmats)',
        '• Quantités (x1 → x99)',
        '• Blocs de texte libres',
        '• Export texte format Discord',
        '• Import depuis un message existant',
    ]):
        p = tr.paragraphs[0] if i == 0 else tr.add_paragraph()
        p.text = line
        p.font.size = Pt(14 if i > 0 else 18)
        p.font.bold = (i == 0)
        p.font.color.rgb = COLORS['text_light'] if i == 0 else COLORS['text_main']
        p.space_after = Pt(6)

    add_footer(slide, '💡  Basculer entre les modes : boutons "🎨 Template" / "📋 MPF TodoList" dans la toolbar')


def create_template_features_slide(prs):
    """Slide 4 — Mode Template : fonctionnalités."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, '🎨  Mode Template — Fonctionnalités')
    add_separator(slide, 1.2)

    bullets = [
        '🎯 Canvas',
        '  Taille fixe 1920×1080 px (proportions Discord parfaites)',
        '  Fond de carte parmi les arrière-plans Foxhole disponibles',
        '',
        '🎯 Sections',
        '  Double-clic sur le canvas → créer une section',
        '  Drag du header → déplacer la section',
        '  Drag du coin → redimensionner',
        '  Titre et couleur personnalisables',
        '',
        '🎯 Icônes',
        '  Plus de 300 items Foxhole disponibles dans la sidebar',
        '  Filtre par catégorie (Armes, Munitions, Véhicules, etc.)',
        '  Drag depuis la sidebar → dépôt dans une section',
        '  Drag à l\'intérieur d\'une section → réordonner',
    ]
    add_bullet_text(slide, bullets, font_size=16, top=1.45, width=6.8)

    add_placeholder_box(
        slide,
        'CAPTURE :\nCanvas avec plusieurs sections\n+ sidebar icônes',
        left=7.6, top=1.45, width=5.2, height=4.75,
    )


def create_template_workflow_slide(prs):
    """Slide 5 — Mode Template : workflow pas à pas."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, '🔄  Mode Template — Workflow')
    add_separator(slide, 1.2)

    # Steps en flow vertical
    steps = [
        ('1', 'Double-clic sur le canvas',      2.2, 1.55),
        ('2', 'Nommez la section + couleur',     2.2, 2.55),
        ('3', 'Drag & drop un icône (sidebar)',  2.2, 3.55),
        ('4', 'Répétez / réorganisez',           2.2, 4.55),
        ('5', 'Ctrl+E → Export PNG',             2.2, 5.55),
    ]
    for num, text, l, t in steps:
        add_step_box(slide, num, text, l, t, width=3.5, height=0.7)
        if num != '5':
            add_arrow(slide, l + 1.55, t + 0.7, direction='down')

    # Flèche droite entre le flow et le résultat
    add_arrow(slide, 5.9, 3.9, direction='right')

    # Box résultat
    result_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.5), Inches(1.55), Inches(6.3), Inches(4.75),
    )
    result_box.fill.solid()
    result_box.fill.fore_color.rgb = COLORS['box_bg']
    result_box.line.color.rgb = COLORS['accent']
    result_box.line.width = Pt(2)

    tf = result_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_top = Pt(14)
    for i, line in enumerate([
        '✅ Résultat',
        '',
        'PNG 1920×1080 prêt à partager',
        'sur Discord / forum Foxhole',
        '',
        'Fond de carte visible sur l\'export',
        'Sections avec titres colorés',
        'Icônes avec quantités',
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(15 if i > 0 else 20)
        p.font.bold = (i == 0)
        p.font.color.rgb = COLORS['accent'] if i == 0 else COLORS['text_main']
        p.space_after = Pt(8)


def create_template_export_slide(prs):
    """Slide 6 — Export PNG."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, '📤  Export PNG — 1920×1080')
    add_separator(slide, 1.2)

    bullets = [
        'Raccourci clavier : Ctrl + E',
        '',
        'Rendu fidèle grâce à html2canvas-pro',
        '  Support complet oklab/oklch (Tailwind v4)',
        '  Fond de carte inclus dans l\'export',
        '  Sections et icônes en haute qualité',
        '',
        'Sauvegarde / Chargement JSON',
        '  Ctrl + S → sauvegarde (export JSON)',
        '  Ctrl + O → chargement (import JSON)',
        '  Auto-save dans localStorage à chaque modification',
        '',
        '✨ Format idéal pour les posts Discord et les forums',
        '✨ Compatible avec tous les fonds de carte Foxhole',
    ]
    add_bullet_text(slide, bullets, font_size=16, top=1.45, width=6.5)

    # Tableau raccourcis clavier
    kbd_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.0), Inches(1.45), Inches(5.8), Inches(4.75),
    )
    kbd_box.fill.solid()
    kbd_box.fill.fore_color.rgb = COLORS['box_bg']
    kbd_box.line.color.rgb = COLORS['primary']
    kbd_box.line.width = Pt(1.5)

    tf = kbd_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_top = Pt(12)
    shortcuts = [
        ('⌨️  Raccourcis clavier', True, 17, COLORS['text_light']),
        ('', False, 10, COLORS['text_muted']),
        ('Ctrl + S   →  Sauvegarder (JSON)', False, 14, COLORS['text_main']),
        ('Ctrl + O   →  Charger un fichier', False, 14, COLORS['text_main']),
        ('Ctrl + E   →  Exporter en PNG', False, 14, COLORS['text_main']),
        ('?              →  Aide / raccourcis', False, 14, COLORS['text_main']),
        ('Echap         →  Fermer les modales', False, 14, COLORS['text_main']),
        ('', False, 10, COLORS['text_muted']),
        ('🖱️  Souris', True, 17, COLORS['text_light']),
        ('', False, 10, COLORS['text_muted']),
        ('Double-clic canvas   →  Nouvelle section', False, 13, COLORS['text_main']),
        ('Drag header          →  Déplacer section', False, 13, COLORS['text_main']),
        ('Drag coin            →  Redimensionner', False, 13, COLORS['text_main']),
        ('Clic icône placé     →  Éditer quantité', False, 13, COLORS['text_main']),
    ]
    for i, (text, bold, size, color) in enumerate(shortcuts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(4)


def create_todolist_features_slide(prs):
    """Slide 7 — Mode MPF Todolist : fonctionnalités."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, '📋  MPF Todolist — Fonctionnalités')
    add_separator(slide, 1.2)

    bullets = [
        '🎯 Catalogue MPF complet',
        '  Toutes les catégories : Small Arms, Heavy Arms, Munitions,',
        '  Véhicules, Navires, Uniformes, Fournitures',
        '  Recherche par nom + filtre par faction (Warden / Colonial)',
        '',
        '🎯 Calcul automatique des coûts',
        '  Bmats / Rmats / Emats / HEmats calculés selon la quantité',
        '  Coût total affiché en temps réel dans le preview',
        '',
        '🎯 Blocs de texte libres',
        '  Ajouter des notes au-dessus, en bas ou par catégorie',
        '  Supporte le markdown Discord (gras, italique, souligné)',
        '',
        '🎯 Import depuis Discord',
        '  Coller un message Discord existant → parsing automatique',
    ]
    add_bullet_text(slide, bullets, font_size=15, top=1.45, width=6.5)

    add_placeholder_box(
        slide,
        'CAPTURE :\nInterface Todolist\n(items + sidebar + preview)',
        left=7.0, top=1.45, width=5.8, height=4.75,
    )


def create_todolist_format_slide(prs):
    """Slide 8 — Format de sortie Discord."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, '📤  Format de sortie — Discord')
    add_separator(slide, 1.2)

    # Colonne explication
    bullets = [
        'Génère du texte prêt à coller dans Discord',
        '',
        'Structure du message :',
        '  __**TITRE DD/MM**__',
        '  (date auto si activée)',
        '',
        '  **__Catégorie__**',
        '  🇦・Item Name – X Bmats (x2)',
        '  🇧・Item Name – X Rmats',
        '  ...',
        '',
        '✨ Lettres indicatrices régionales (🇦 🇧 🇨…)',
        '✨ Coûts MPF calculés et affichés',
        '✨ Multiplicateur (x2, x5…) inclus',
        '✨ Blocs de texte libres positionnés',
    ]
    add_bullet_text(slide, bullets, font_size=15, top=1.45, width=6.2)

    # Box exemple de sortie
    example_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.45), Inches(6.0), Inches(4.75),
    )
    example_box.fill.solid()
    example_box.fill.fore_color.rgb = COLORS['tag_bg']
    example_box.line.color.rgb = RGBColor(88, 101, 242)
    example_box.line.width = Pt(2)

    tf = example_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(16)
    tf.margin_top = Pt(14)
    example_lines = [
        ('Exemple de sortie :', True, 16, COLORS['accent_gold']),
        ('', False, 8,  COLORS['text_muted']),
        ('__**LOGI TODOLIST 14/05**__', False, 14, COLORS['text_light']),
        ('', False, 8,  COLORS['text_muted']),
        ('**__Small Arms__**', False, 13, COLORS['accent']),
        ('🇦・Argenti r.II – 100 Bmats', False, 13, COLORS['text_main']),
        ('🇧・Catena r.II – 100 Bmats (x2)', False, 13, COLORS['text_main']),
        ('', False, 8,  COLORS['text_muted']),
        ('**__Vehicles__**', False, 13, COLORS['accent']),
        ('🇦・Dunne Loadlugger – 250 Bmats', False, 13, COLORS['text_main']),
        ('🇧・BMS Mineseeker – 150 Bmats', False, 13, COLORS['text_main']),
        ('  – 25 Rmats', False, 12, COLORS['text_muted']),
    ]
    for i, (text, bold, size, color) in enumerate(example_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(3)


def create_todolist_import_slide(prs):
    """Slide 9 — Import / Preview en temps réel."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, '🔄  Import & Preview en temps réel')
    add_separator(slide, 1.2)

    bullets = [
        '📥 Import depuis Discord',
        '  Bouton "Charger" → coller un message Discord existant',
        '  Parsing automatique du format __**Titre**__ + items',
        '  Lettres 🇦 🇧 🇨… ou A- B- C- reconnues',
        '  Quantités (x5) et coûts préservés',
        '',
        '👁️ Preview en temps réel',
        '  Affichage côte à côte éditeur / rendu final',
        '  Emojis Discord renderisés (⚠️ ✅ ❌ ⏰…)',
        '  Markdown Discord simulé (gras, italique, souligné)',
        '  Bouton copie en un clic',
        '',
        '💾 Auto-save',
        '  La todolist est sauvegardée automatiquement dans localStorage',
        '  Rechargement de page → retrouve l\'état précédent',
    ]
    add_bullet_text(slide, bullets, font_size=15, top=1.45, width=6.5)

    add_placeholder_box(
        slide,
        'CAPTURE :\nModal d\'import\n+\nPreview Discord simulé',
        left=7.0, top=1.45, width=5.8, height=4.75,
    )


def create_tech_stack_slide(prs):
    """Slide 10 — Stack technique."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, '🛠️  Stack technique')
    add_separator(slide, 1.2)

    # Colonne gauche
    left_items = [
        ('TypeScript',       'Typage statique, maintenabilité'),
        ('Vite',             'Build ultra-rapide, HMR'),
        ('Tailwind CSS v4',  'Utility-first, styles modernes'),
        ('interact.js',      'Drag/resize des sections'),
        ('SortableJS',       'Réordonnancement des icônes'),
    ]
    right_items = [
        ('html2canvas-pro',  'Export PNG (support oklch/oklab)'),
        ('JSON local',       'Sauvegarde templates + todolists'),
        ('Vite + GitHub Pages', 'CI/CD automatique'),
        ('Jest + ts-jest',   'Tests unitaires TypeScript'),
        ('CC BY-NC 4.0',     'Licence open-source'),
    ]

    for col_x, items in [(0.5, left_items), (6.9, right_items)]:
        for row, (name, desc) in enumerate(items):
            top = 1.55 + row * 0.95

            tag_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(col_x), Inches(top), Inches(2.4), Inches(0.5),
            )
            tag_shape.fill.solid()
            tag_shape.fill.fore_color.rgb = COLORS['primary']
            tag_shape.line.fill.background()
            tf_tag = tag_shape.text_frame
            p_tag = tf_tag.paragraphs[0]
            p_tag.text = name
            p_tag.font.size = Pt(13)
            p_tag.font.bold = True
            p_tag.font.color.rgb = COLORS['text_light']
            p_tag.alignment = PP_ALIGN.CENTER
            tf_tag.anchor = MSO_ANCHOR.MIDDLE

            desc_box = slide.shapes.add_textbox(
                Inches(col_x + 2.6), Inches(top + 0.05),
                Inches(3.8), Inches(0.45),
            )
            tf_desc = desc_box.text_frame
            p_desc = tf_desc.paragraphs[0]
            p_desc.text = desc
            p_desc.font.size = Pt(13)
            p_desc.font.color.rgb = COLORS['text_main']

    add_footer(slide, '📦  npm install  →  npm run dev  |  Tests : npm test')


def create_end_slide(prs):
    """Slide 11 — Fin / Questions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Titre centré
    title = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.333), Inches(1.4))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = '🎮  Prêt à l\'emploi !'
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.CENTER

    add_separator(slide, 3.55, width=12.333)

    sub = slide.shapes.add_textbox(Inches(0.5), Inches(3.75), Inches(12.333), Inches(0.8))
    tf2 = sub.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = 'Questions / retours ? Ouvrez une issue sur GitHub'
    p2.font.size = Pt(24)
    p2.font.italic = True
    p2.font.color.rgb = COLORS['text_muted']
    p2.alignment = PP_ALIGN.CENTER

    # Badges bas de page
    add_tag(slide, '📄  CC BY-NC 4.0', 2.0, 5.0, width=2.8, height=0.5,
            bg=RGBColor(60, 80, 60), border=RGBColor(60, 120, 80), font_size=14)
    add_tag(slide, '🌐  GitHub Pages', 5.2, 5.0, width=2.8, height=0.5,
            bg=COLORS['primary'], border=COLORS['primary'], font_size=14)
    add_tag(slide, '🦊  Foxhole by Siege Camp', 8.4, 5.0, width=3.4, height=0.5,
            bg=COLORS['accent_orange'], border=COLORS['accent_orange'], font_size=14)

    note = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.6))
    tf3 = note.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = 'Made with ❤️ for the 11eRC-FL  —  v2.0'
    p3.font.size = Pt(14)
    p3.font.italic = True
    p3.font.color.rgb = COLORS['text_muted']
    p3.alignment = PP_ALIGN.CENTER


# ─────────────────────────────────────────────
#  Point d'entrée
# ─────────────────────────────────────────────

def create_presentation():
    """Crée la présentation PowerPoint complète."""
    prs = Presentation()
    prs.slide_width  = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)

    create_title_slide(prs)
    create_overview_slide(prs)
    create_modes_slide(prs)
    create_template_features_slide(prs)
    create_template_workflow_slide(prs)
    create_template_export_slide(prs)
    create_todolist_features_slide(prs)
    create_todolist_format_slide(prs)
    create_todolist_import_slide(prs)
    create_tech_stack_slide(prs)
    create_end_slide(prs)

    # Fond sur toutes les slides
    for s in prs.slides:
        add_background(s, prs)

    output_path = os.path.join(SCRIPT_DIR, '11e_TemplateBuilder_Presentation.pptx')
    prs.save(output_path)
    print(f'✅  Présentation générée : {output_path}')
    return output_path


if __name__ == '__main__':
    create_presentation()
